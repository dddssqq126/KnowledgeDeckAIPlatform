"""Demo: KB ingestion (parse -> chunk -> dual-embed -> Qdrant upsert).

Standalone re-implementation of the production pipeline. Mirrors:
  - backend/app/features/rag/services/document_parser.py
  - backend/app/features/rag/services/text_splitter.py
  - backend/app/features/rag/services/sparse_embed.py
  - backend/app/features/rag/services/qdrant_store.py
  - backend/app/features/rag/services/ingestion.py
as of commit a9ad2d5. If those files diverge, sync this file by hand.

Usage:
  python 01_kb_ingest.py path/to/file.pdf
  python 01_kb_ingest.py path/to/file.txt --filename my_notes.txt --file-id 7
  python 01_kb_ingest.py path/to/file.pdf --cleanup-first

Supported extensions: txt, md, pdf, docx, pptx, py, html, css, cs.
"""
from __future__ import annotations

import argparse
import io
import os
import sys
import uuid

import httpx
from docx import Document as DocxDocument
from fastembed import SparseTextEmbedding
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pypdf import PdfReader
from qdrant_client import QdrantClient
from qdrant_client.http import models as qm

from _common import (
    DEMO_KB_ID,
    DEMO_USER_ID,
    DENSE_VEC,
    EMBED_API_KEY,
    EMBED_DIM,
    EMBED_MODEL,
    EMBED_URL,
    QDRANT_COLLECTION,
    QDRANT_URL,
    SPARSE_VEC,
    cleanup_demo_vectors,
)

CHUNK_CHARS = 1200
CHUNK_OVERLAP = 150
TEXT_EXTENSIONS = ("txt", "md", "py", "html", "css", "cs")


# ---------------------------------------------------------------- 1. Parse
# Per-format parsers. PDF/PPTX emit one segment per page/slide so we keep
# page numbers in the payload; everything else collapses to a single
# segment with page_number=None.

def parse_file(extension: str, data: bytes) -> list[tuple[str, int | None]]:
    if extension in TEXT_EXTENSIONS:
        return [(data.decode("utf-8", errors="replace"), None)]

    if extension == "pdf":
        out: list[tuple[str, int | None]] = []
        for i, page in enumerate(PdfReader(io.BytesIO(data)).pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                out.append((text, i))
        return out

    if extension == "docx":
        # Word has no exposed page concept (page breaks are layout decisions),
        # so we flatten paragraphs + table cells into one segment.
        doc = DocxDocument(io.BytesIO(data))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
        return [(text, None)] if text.strip() else []

    if extension == "pptx":
        out = []
        for i, slide in enumerate(Presentation(io.BytesIO(data)).slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            text = "\n".join(parts)
            if text.strip():
                out.append((text, i))
        return out

    raise ValueError(f"unsupported extension: {extension}")


# ---------------------------------------------------------------- 2. Chunk
# RecursiveCharacterTextSplitter walks separators in priority order so
# breakpoints land on natural boundaries (paragraph -> newline -> sentence
# -> word -> char) rather than mid-token.

def split_into_chunks(text: str) -> list[str]:
    text = text.strip()
    if not text:
        return []
    if len(text) <= CHUNK_CHARS:
        return [text]
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_CHARS,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", "? ", "! ", "; ", ", ", " ", ""],
        length_function=len,
        is_separator_regex=False,
    )
    return [c for c in (s.strip() for s in splitter.split_text(text)) if c]


# ---------------------------------------------------------------- 3. Dense embed
# vLLM serves an OpenAI-compatible /v1/embeddings endpoint; bge-m3 returns
# 1024-d vectors. We batch all chunks into a single request.

def dense_embed(texts: list[str]) -> list[list[float]]:
    payload = {"model": EMBED_MODEL, "input": texts}
    headers = {"Authorization": f"Bearer {EMBED_API_KEY}"}
    with httpx.Client(timeout=60.0) as client:
        r = client.post(
            f"{EMBED_URL.rstrip('/')}/embeddings",
            json=payload,
            headers=headers,
        )
        r.raise_for_status()
        return [item["embedding"] for item in r.json()["data"]]


# ---------------------------------------------------------------- 4. Sparse embed
# fastembed Qdrant/bm25 ships pre-computed IDF tables, so it's
# corpus-free. The model loads on first .embed() call (~tens of MB).

def sparse_embed(texts: list[str]) -> list[tuple[list[int], list[float]]]:
    model = SparseTextEmbedding(model_name="Qdrant/bm25")
    out: list[tuple[list[int], list[float]]] = []
    for v in model.embed(list(texts)):
        out.append(
            (
                [int(i) for i in v.indices.tolist()],
                [float(x) for x in v.values.tolist()],
            )
        )
    return out


# ---------------------------------------------------------------- 5. Qdrant write
# Hybrid named-vector schema: each point carries a `dense` (cosine) +
# `sparse` (BM25) vector. Per-user filtering is enforced via payload
# filters at query time, not via separate collections.

def ensure_collection(client: QdrantClient) -> None:
    if client.collection_exists(QDRANT_COLLECTION):
        return
    client.create_collection(
        collection_name=QDRANT_COLLECTION,
        vectors_config={
            DENSE_VEC: qm.VectorParams(size=EMBED_DIM, distance=qm.Distance.COSINE),
        },
        sparse_vectors_config={
            SPARSE_VEC: qm.SparseVectorParams(index=qm.SparseIndexParams(on_disk=False)),
        },
    )
    for field in ("user_id", "kb_id", "file_id"):
        client.create_payload_index(
            collection_name=QDRANT_COLLECTION,
            field_name=field,
            field_schema=qm.PayloadSchemaType.INTEGER,
        )


def upsert_points(
    client: QdrantClient,
    *,
    file_id: int,
    filename: str,
    chunks: list[dict],
    dense: list[list[float]],
    sparse: list[tuple[list[int], list[float]]],
) -> None:
    points = [
        qm.PointStruct(
            id=str(uuid.uuid4()),
            vector={
                DENSE_VEC: d,
                SPARSE_VEC: qm.SparseVector(indices=s[0], values=s[1]),
            },
            payload={
                "user_id": DEMO_USER_ID,
                "kb_id": DEMO_KB_ID,
                "file_id": file_id,
                "filename": filename,
                "text": ch["text"],
                "page_number": ch.get("page_number"),
                "chunk_index": ch["chunk_index"],
            },
        )
        for ch, d, s in zip(chunks, dense, sparse, strict=True)
    ]
    client.upsert(collection_name=QDRANT_COLLECTION, points=points)


# ---------------------------------------------------------------- 6. CLI driver

def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("path", help="path to file to ingest")
    ap.add_argument("--filename", help="display filename in citations (default: basename)")
    ap.add_argument("--file-id", type=int, default=1, help="file_id payload (default: 1)")
    ap.add_argument(
        "--cleanup-first",
        action="store_true",
        help="wipe existing DEMO_USER_ID points before ingest",
    )
    args = ap.parse_args()

    if args.cleanup_first:
        n = cleanup_demo_vectors()
        print(f"[cleanup] dropped {n} prior demo points")

    if not os.path.exists(args.path):
        print(f"file not found: {args.path}", file=sys.stderr)
        return 2

    extension = args.path.rsplit(".", 1)[-1].lower()
    filename = args.filename or os.path.basename(args.path)
    with open(args.path, "rb") as f:
        data = f.read()

    print(f"[parse]   extension={extension} filename={filename} bytes={len(data)}")
    segments = parse_file(extension, data)
    if not segments:
        print("no extractable text", file=sys.stderr)
        return 1
    print(f"[parse]   {len(segments)} segment(s)")

    chunks: list[dict] = []
    for text, page in segments:
        for piece in split_into_chunks(text):
            chunks.append(
                {"text": piece, "page_number": page, "chunk_index": len(chunks)}
            )
    print(f"[chunk]   {len(chunks)} chunk(s)  (size={CHUNK_CHARS}, overlap={CHUNK_OVERLAP})")
    if not chunks:
        return 1

    texts = [c["text"] for c in chunks]
    print(f"[dense]   embedding {len(texts)} chunk(s) via {EMBED_URL} ({EMBED_MODEL}, {EMBED_DIM}d)")
    dense_vecs = dense_embed(texts)
    print(f"[sparse]  embedding {len(texts)} chunk(s) via fastembed Qdrant/bm25")
    sparse_vecs = sparse_embed(texts)

    print(f"[qdrant]  upserting to {QDRANT_URL}/{QDRANT_COLLECTION}")
    client = QdrantClient(url=QDRANT_URL)
    ensure_collection(client)
    upsert_points(
        client,
        file_id=args.file_id,
        filename=filename,
        chunks=chunks,
        dense=dense_vecs,
        sparse=sparse_vecs,
    )
    print(
        f"[done]    {len(chunks)} points written under "
        f"user_id={DEMO_USER_ID} kb_id={DEMO_KB_ID} file_id={args.file_id}"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
